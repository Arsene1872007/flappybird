"""
Merge two presentations:
  - user's 'flappy bird.pptx' (all 7 slides)
  - my FlappyBird_Presentation.pptx (slides 2-10, skip first title slide)

Output: FlappyBird_Final.pptx at 13.33 x 7.5 inches
"""
from pptx import Presentation
from pptx.util import Inches, Emu
from lxml import etree
import copy

SRC_USER = r"C:\Users\shemo\Downloads\flappy bird.pptx"
SRC_MINE = r"C:\Users\shemo\OneDrive\Desktop\flappybird\FlappyBird_Presentation.pptx"
OUT      = r"C:\Users\shemo\OneDrive\Desktop\flappybird\FlappyBird_Final.pptx"

TARGET_W = Inches(13.33)
TARGET_H = Inches(7.5)

# ── helpers ───────────────────────────────────────────────────────────────────

def copy_slide(src_prs, slide_index, dst_prs):
    """
    Deep-copy one slide (by index) from src_prs into dst_prs,
    carrying along all media, themes, and layout references.
    Adapted from standard lxml-based copy pattern for python-pptx.
    """
    template_slide = src_prs.slides[slide_index]

    # We need a blank layout in the destination to host the new slide
    blank_layout = dst_prs.slide_layouts[6]
    new_slide = dst_prs.slides.add_slide(blank_layout)

    # Copy the entire spTree (shapes) from source
    src_sp_tree = template_slide.shapes._spTree
    dst_sp_tree = new_slide.shapes._spTree

    # Remove placeholder shapes that came from blank layout
    for el in list(dst_sp_tree):
        dst_sp_tree.remove(el)

    # Deep-copy every element from the source slide's spTree
    for el in src_sp_tree:
        dst_sp_tree.append(copy.deepcopy(el))

    # Copy slide background (if any)
    src_bg = template_slide.background
    dst_bg = new_slide.background
    src_fill = src_bg.fill
    # Copy the bg element directly
    src_bg_elem = template_slide._element.find(
        './/{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}bg'
    )
    # Use the slide-level bg element
    src_cSld = template_slide._element.find(
        '{http://schemas.openxmlformats.org/presentationml/2006/main}cSld'
    )
    dst_cSld = new_slide._element.find(
        '{http://schemas.openxmlformats.org/presentationml/2006/main}cSld'
    )

    if src_cSld is not None and dst_cSld is not None:
        # Copy bg element from source if present
        src_bg_el = src_cSld.find(
            '{http://schemas.openxmlformats.org/presentationml/2006/main}bg'
        )
        if src_bg_el is not None:
            dst_bg_el = dst_cSld.find(
                '{http://schemas.openxmlformats.org/presentationml/2006/main}bg'
            )
            if dst_bg_el is not None:
                dst_cSld.remove(dst_bg_el)
            dst_cSld.insert(0, copy.deepcopy(src_bg_el))

    # Copy image relationships (pictures will be broken without this)
    for rel in template_slide.part.rels.values():
        if "image" in rel.reltype:
            img_part = rel.target_part
            new_rId = new_slide.part.relate_to(img_part, rel.reltype)
            # Fix up rId references in the new slide XML
            old_rId = rel.rId
            if old_rId != new_rId:
                xml_str = etree.tostring(new_slide._element, encoding="unicode")
                xml_str = xml_str.replace(f'r:embed="{old_rId}"', f'r:embed="{new_rId}"')
                xml_str = xml_str.replace(f'r:id="{old_rId}"', f'r:id="{new_rId}"')
                new_slide._element[:] = []
                new_elem = etree.fromstring(xml_str)
                for child in new_elem:
                    new_slide._element.append(child)

    # Copy audio/video relationships
    for rel in template_slide.part.rels.values():
        if "audio" in rel.reltype or "video" in rel.reltype:
            try:
                media_part = rel.target_part
                new_slide.part.relate_to(media_part, rel.reltype)
            except Exception:
                pass

    return new_slide


def scale_slide_shapes(slide, src_w, src_h, dst_w, dst_h):
    """Rescale all shape positions/sizes to fit the new slide dimensions."""
    sx = dst_w / src_w
    sy = dst_h / src_h
    for shape in slide.shapes:
        shape.left   = int(shape.left   * sx)
        shape.top    = int(shape.top    * sy)
        shape.width  = int(shape.width  * sx)
        shape.height = int(shape.height * sy)


# ── build destination ─────────────────────────────────────────────────────────

dst = Presentation()
dst.slide_width  = TARGET_W
dst.slide_height = TARGET_H

user_prs = Presentation(SRC_USER)
mine_prs = Presentation(SRC_MINE)

USER_W = user_prs.slide_width
USER_H = user_prs.slide_height
MINE_W = mine_prs.slide_width
MINE_H = mine_prs.slide_height

print(f"User slides: {len(user_prs.slides)}  size: {USER_W.inches:.2f}x{USER_H.inches:.2f}")
print(f"My slides:   {len(mine_prs.slides)}  size: {MINE_W.inches:.2f}x{MINE_H.inches:.2f}")

# Add all 7 user slides
for i in range(len(user_prs.slides)):
    s = copy_slide(user_prs, i, dst)
    if USER_W != TARGET_W or USER_H != TARGET_H:
        scale_slide_shapes(s, USER_W, USER_H, TARGET_W, TARGET_H)
    print(f"  Added user slide {i+1}")

# Add my slides 2-10 (skip index 0 = title slide)
for i in range(1, len(mine_prs.slides)):
    s = copy_slide(mine_prs, i, dst)
    # My slides are already TARGET_W x TARGET_H — no rescaling needed
    print(f"  Added my slide {i+1}")

dst.save(OUT)
print(f"\nSaved: {OUT}")
print(f"Total slides: {len(dst.slides)}")
