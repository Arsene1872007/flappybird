"""Build Flappy Bird presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
DARK_BG   = RGBColor(0x0D, 0x1B, 0x2A)   # deep night sky
MID_BG    = RGBColor(0x1B, 0x35, 0x53)   # medium blue
SKY_BLUE  = RGBColor(0x4E, 0xC9, 0xF0)   # bright cyan
GOLD      = RGBColor(0xF5, 0xA6, 0x23)   # score gold
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TXT = RGBColor(0xC8, 0xE6, 0xF7)   # soft light text
GREEN     = RGBColor(0x4C, 0xAF, 0x50)   # pipe green
ORANGE    = RGBColor(0xFF, 0x7B, 0x25)   # warm accent

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ── Helpers ───────────────────────────────────────────────────────────────────

def blank_slide():
    layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(layout)

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def section_label(slide, label, color=GOLD):
    rect(slide, 0, 0, 13.33, 0.55, MID_BG)
    txt(slide, label, 0.35, 0.05, 12, 0.5, size=14, bold=True, color=color,
        align=PP_ALIGN.LEFT)

def title_area(slide, title, subtitle=None):
    txt(slide, title, 0.5, 1.4, 12.33, 1.5, size=52, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        txt(slide, subtitle, 0.5, 3.1, 12.33, 0.8, size=22, bold=False,
            color=LIGHT_TXT, align=PP_ALIGN.CENTER)

def bullet_rows(slide, items, x, y, w, icon_color=SKY_BLUE, text_color=WHITE,
                size=17, row_h=0.58):
    for i, (icon, text) in enumerate(items):
        iy = y + i * row_h
        # icon circle
        circle = slide.shapes.add_shape(9, Inches(x), Inches(iy + 0.03),
                                         Inches(0.35), Inches(0.35))
        circle.fill.solid(); circle.fill.fore_color.rgb = icon_color
        circle.line.fill.background()
        txt(slide, icon, x + 0.01, iy, 0.35, 0.38, size=13, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        txt(slide, text, x + 0.50, iy, w - 0.55, 0.42, size=size,
            color=text_color, align=PP_ALIGN.LEFT)

def card(slide, x, y, w, h, fill=MID_BG, radius=False):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(); bg(s, DARK_BG)

# decorative pipe shapes
for gx, gy, gw, gh in [(0, 0, 1.2, 2.8), (11.8, 4.5, 1.53, 3.05)]:
    rect(s, gx, gy, gw, gh, GREEN)

# star dots
for dx, dy in [(2,0.5),(5,0.2),(8,0.8),(10.5,0.3),(1.5,5),(4,5.8),(7.2,5.5),(11,5.2)]:
    dot = s.shapes.add_shape(9, Inches(dx), Inches(dy), Inches(0.08), Inches(0.08))
    dot.fill.solid(); dot.fill.fore_color.rgb = WHITE; dot.line.fill.background()

txt(s, "🐦 FLAPPY BIRD", 0.5, 1.5, 12.33, 2.0, size=64, bold=True,
    color=GOLD, align=PP_ALIGN.CENTER)
txt(s, "A Python + Pygame Side-Scrolling Game", 0.5, 3.35, 12.33, 0.8,
    size=22, color=LIGHT_TXT, align=PP_ALIGN.CENTER)
txt(s, "10-Minute Feature Walk-Through", 0.5, 4.2, 12.33, 0.6,
    size=16, italic=True, color=SKY_BLUE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — What is it?
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(); bg(s, DARK_BG)
section_label(s, "OVERVIEW")

# big tagline
txt(s, "Dodge pipes. Beat your score. Use your voice.", 0.5, 0.7, 12.33, 1.2,
    size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 3 stat cards
for i, (num, label, clr) in enumerate([
    ("Python", "Language", SKY_BLUE),
    ("2 Modes", "Normal + Voice", GOLD),
    ("Endless", "Side-Scrolling", GREEN),
]):
    cx = 0.6 + i * 4.25
    card(s, cx, 2.0, 3.8, 2.2, fill=MID_BG)
    txt(s, num, cx + 0.1, 2.15, 3.6, 1.0, size=34, bold=True, color=clr,
        align=PP_ALIGN.CENTER)
    txt(s, label, cx + 0.1, 3.05, 3.6, 0.7, size=15, color=LIGHT_TXT,
        align=PP_ALIGN.CENTER)

txt(s, "Built with Pygame · sounddevice · numpy  |  All assets hot-swappable",
    0.5, 4.5, 12.33, 0.5, size=13, italic=True, color=LIGHT_TXT,
    align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Core Gameplay
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(); bg(s, DARK_BG)
section_label(s, "CORE GAMEPLAY")

txt(s, "Physics, Pipes & Points", 0.5, 0.65, 12.33, 0.9, size=34, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER)

items = [
    ("↓", "Gravity pulls the bird down continuously — always falling"),
    ("tap", "Flap to push upward; release and gravity takes over"),
    ("🪧", "Pipes spawn at random heights — endless side-scroll"),
    ("+1", "Score increases by 1 for every pipe you clear"),
    ("💥", "Collide with a pipe, the ground, or the ceiling → game over"),
]
bullet_rows(s, items, 1.2, 1.7, 11.0, icon_color=GREEN, row_h=0.72)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Controls: Two Modes
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(); bg(s, DARK_BG)
section_label(s, "CONTROLS")

txt(s, "Two Ways to Play", 0.5, 0.65, 12.33, 0.9, size=34, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER)

# Normal mode card
card(s, 0.5, 1.75, 5.8, 4.8, fill=MID_BG)
rect(s, 0.5, 1.75, 5.8, 0.5, SKY_BLUE)
txt(s, "NORMAL MODE", 0.55, 1.75, 5.7, 0.5, size=16, bold=True,
    color=DARK_BG, align=PP_ALIGN.CENTER)
normal_items = [
    ("⌨", "Press SPACE to flap"),
    ("🖱", "Left-click to flap"),
    ("🏁", "Available immediately — no setup"),
    ("✅", "Works on any machine"),
]
bullet_rows(s, normal_items, 0.75, 2.45, 5.2, icon_color=SKY_BLUE, row_h=0.78)

# Voice mode card
card(s, 6.9, 1.75, 5.93, 4.8, fill=MID_BG)
rect(s, 6.9, 1.75, 5.93, 0.5, ORANGE)
txt(s, "VOICE MODE", 6.95, 1.75, 5.8, 0.5, size=16, bold=True,
    color=DARK_BG, align=PP_ALIGN.CENTER)
voice_items = [
    ("🎤", "Speak or make noise to flap"),
    ("📊", "Louder sound = stronger lift"),
    ("📈", "Live mic-level meter on screen"),
    ("⚡", "6 ms low-latency audio engine"),
]
bullet_rows(s, voice_items, 7.15, 2.45, 5.4, icon_color=ORANGE, row_h=0.78)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Voice Mode Deep Dive
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(); bg(s, DARK_BG)
section_label(s, "VOICE MODE — DEEP DIVE")

txt(s, "How Your Voice Controls the Bird", 0.5, 0.65, 12.33, 0.9, size=32,
    bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# flow steps
steps = [
    ("1", "Mic captures audio via sounddevice", SKY_BLUE),
    ("2", "numpy computes RMS volume level",    GOLD),
    ("3", "Level compared to threshold",         ORANGE),
    ("4", "Above threshold → upward flap force", GREEN),
    ("5", "Louder sound → stronger push",        SKY_BLUE),
]
for i, (num, desc, clr) in enumerate(steps):
    bx = 0.5 + i * 2.5
    card(s, bx, 1.8, 2.25, 1.5, fill=MID_BG)
    txt(s, num, bx + 0.1, 1.9, 2.05, 0.7, size=36, bold=True, color=clr,
        align=PP_ALIGN.CENTER)
    txt(s, desc, bx + 0.1, 2.55, 2.05, 0.65, size=12, color=LIGHT_TXT,
        align=PP_ALIGN.CENTER)
    if i < 4:
        txt(s, "→", bx + 2.28, 2.1, 0.22, 0.6, size=22, bold=True,
            color=GOLD, align=PP_ALIGN.CENTER)

# extra feature boxes
for bx, title, body, clr in [
    (0.5, "Live Volume Meter",
     "A real-time bar on screen shows mic level vs. threshold so you can calibrate your voice", SKY_BLUE),
    (4.6, "Graceful Fallback",
     "If sounddevice / numpy aren't installed, Voice Mode is disabled with a clear message", ORANGE),
    (8.7, "Mode Badge",
     "A corner badge always tells you which control mode is currently active", GOLD),
]:
    card(s, bx, 3.6, 3.85, 2.5, fill=MID_BG)
    rect(s, bx, 3.6, 3.85, 0.45, clr)
    txt(s, title, bx + 0.1, 3.62, 3.65, 0.42, size=14, bold=True,
        color=DARK_BG, align=PP_ALIGN.CENTER)
    txt(s, body, bx + 0.15, 4.12, 3.55, 1.85, size=13, color=LIGHT_TXT,
        wrap=True)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Scoring & Progression
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(); bg(s, DARK_BG)
section_label(s, "SCORING & PROGRESSION")

txt(s, "Keep Your Best Score — Forever", 0.5, 0.65, 12.33, 0.9, size=32,
    bold=True, color=WHITE, align=PP_ALIGN.CENTER)

items = [
    ("💾", "Best score saved to disk — survives app restarts"),
    ("📌", "Best score displayed on-screen throughout every game"),
    ("🏆", "Game-over screen shows current vs. personal best"),
    ("🥇", "Personal bests highlighted in gold — instant recognition"),
    ("🎉", "Record fanfare sound plays every 10 points"),
    ("🔄", "Switch modes from the game-over screen — no full restart"),
]
bullet_rows(s, items, 1.0, 1.7, 11.3, icon_color=GOLD, row_h=0.73)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Audio & Visuals
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(); bg(s, DARK_BG)
section_label(s, "AUDIO & VISUALS")

txt(s, "Every Action Has Feedback", 0.5, 0.65, 12.33, 0.9, size=32,
    bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Left — Audio
card(s, 0.5, 1.7, 5.9, 5.0, fill=MID_BG)
rect(s, 0.5, 1.7, 5.9, 0.48, SKY_BLUE)
txt(s, "🔊  AUDIO", 0.55, 1.72, 5.8, 0.45, size=16, bold=True,
    color=DARK_BG, align=PP_ALIGN.CENTER)
audio_items = [
    ("♪", "Flap sound on every tap or voice trigger"),
    ("🗣", "Two alternating voices per pipe cleared"),
    ("💀", "Losing sound on death"),
    ("👋", "Welcome sound on title screen"),
    ("⚡", "6 ms low-latency audio buffer"),
]
bullet_rows(s, audio_items, 0.75, 2.35, 5.45, icon_color=SKY_BLUE, row_h=0.73)

# Right — Visuals
card(s, 6.9, 1.7, 5.93, 5.0, fill=MID_BG)
rect(s, 6.9, 1.7, 5.93, 0.48, GOLD)
txt(s, "🎨  VISUALS", 6.95, 1.72, 5.8, 0.45, size=16, bold=True,
    color=DARK_BG, align=PP_ALIGN.CENTER)
vis_items = [
    ("🖼", "Title & game-over image (auto text fallback)"),
    ("🐦", "Animated bird tilt — rotates with direction"),
    ("🌫", "Semi-transparent game-over overlay"),
    ("📟", "Live mic-volume meter in Voice Mode"),
    ("🏷", "Mode badge showing current control type"),
]
bullet_rows(s, vis_items, 7.15, 2.35, 5.5, icon_color=GOLD, row_h=0.73)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Tech Stack
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(); bg(s, DARK_BG)
section_label(s, "TECH STACK")

txt(s, "Under the Hood", 0.5, 0.65, 12.33, 0.9, size=34, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER)

stack = [
    ("Python",        "Core language — readable, cross-platform",            SKY_BLUE),
    ("Pygame",        "Game loop, rendering, collision, sprite animation",     GREEN),
    ("sounddevice",   "Real-time microphone capture at 6 ms latency",          ORANGE),
    ("numpy",         "Fast RMS volume computation on audio buffers",          GOLD),
    ("File I/O",      "Best score persisted to disk across sessions",          SKY_BLUE),
]
for i, (lib, desc, clr) in enumerate(stack):
    bx = 0.5 + (i % 3) * 4.25
    by = 1.75 + (i // 3) * 2.3
    card(s, bx, by, 3.9, 2.0, fill=MID_BG)
    rect(s, bx, by, 3.9, 0.48, clr)
    txt(s, lib, bx + 0.1, by + 0.04, 3.7, 0.42, size=18, bold=True,
        color=DARK_BG, align=PP_ALIGN.CENTER)
    txt(s, desc, bx + 0.15, by + 0.6, 3.6, 1.2, size=13, color=LIGHT_TXT,
        wrap=True, align=PP_ALIGN.CENTER)

txt(s, "Assets (images & sounds) are hot-swappable — just replace files in img/ or sound/",
    0.5, 6.9, 12.33, 0.5, size=13, italic=True, color=LIGHT_TXT,
    align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Live Demo
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(); bg(s, DARK_BG)

for gx, gy, gw, gh in [(0, 0, 1.4, 3.0), (11.93, 4.2, 1.4, 3.3)]:
    rect(s, gx, gy, gw, gh, GREEN)

txt(s, "LIVE DEMO", 0.5, 1.5, 12.33, 1.8, size=72, bold=True,
    color=GOLD, align=PP_ALIGN.CENTER)
txt(s, "Let's play — Normal Mode first, then Voice Mode", 0.5, 3.45, 12.33,
    0.8, size=22, color=LIGHT_TXT, align=PP_ALIGN.CENTER)

things = ["Space / click to flap", "Avoid the pipes", "Try to beat the best score"]
for i, t in enumerate(things):
    tx = 1.5 + i * 3.6
    card(s, tx, 4.5, 3.2, 1.0, fill=MID_BG)
    txt(s, t, tx + 0.1, 4.6, 3.0, 0.7, size=14, color=WHITE,
        align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Thank You / Q&A
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(); bg(s, DARK_BG)

for gx, gy, gw, gh in [(0, 3.5, 1.2, 4.0), (12.13, 0, 1.2, 3.0)]:
    rect(s, gx, gy, gw, gh, GREEN)

txt(s, "Thanks for Watching!", 0.5, 1.3, 12.33, 1.5, size=52, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Questions?", 0.5, 3.0, 12.33, 1.0, size=38, bold=True,
    color=GOLD, align=PP_ALIGN.CENTER)

recap = [
    "Gravity physics + pipe obstacles + collision detection",
    "Normal Mode (Space/click) & Voice Mode (microphone)",
    "Persistent best score + record fanfare every 10 pts",
    "Animated sprites · rich audio · hot-swappable assets",
]
for i, line in enumerate(recap):
    txt(s, f"• {line}", 2.0, 4.25 + i * 0.45, 9.33, 0.42, size=14,
        color=LIGHT_TXT, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
out = r"C:\Users\shemo\OneDrive\Desktop\flappybird\FlappyBird_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
